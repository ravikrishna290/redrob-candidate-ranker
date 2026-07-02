import collections
import collections.abc
from pptx import Presentation

def fill_slide_content():
    # Load presentation
    prs = Presentation(r"d:\Redrob AI hackathon\Idea Submission Template _ Redrob.pptx")
    
    # Define mapping from slide prompt substrings to actual answers
    replacements = {
        # Slide 1
        "Team Name :": "Team Name : [Insert Team Name]",
        "Team Leader Name :": "Team Leader Name : [Insert Team Leader Name]",
        "Problem Statement :": "Problem Statement : High-Performance AI Candidate Discovery and Ranking Engine under 5-minute CPU constraint",
        
        # Slide 2
        "What is your proposed solution?": (
            "Our solution is a high-performance, two-stage AI Candidate Discovery and Ranking Engine. It performs dense semantic retrieval followed by a Cross-Encoder reranking model, integrated with custom recruiter-style scoring rules and hard-coded trap detection to filter out fake profiles.\n\n"
            "Key Differentiators:\n"
            "1. Two-Stage Pipeline: Combines bi-encoder speed with cross-encoder deep self-attention.\n"
            "2. Anti-Honeypot Filter: Hard validation of career timeline anomalies and skill inflation to drop honeypots.\n"
            "3. Recruiter Heuristics: Uses notice periods, job stability, and career progression markers instead of raw keyword similarity."
        ),
        
        # Slide 3
        "What are the key requirements extracted from the JD?": (
            "Key Requirements Extracted:\n"
            "- Role: Senior AI Engineer / Founding Team\n"
            "- Core Skills: PyTorch, embeddings, vector databases (FAISS, Pinecone, Qdrant, Milvus, OpenSearch), LLMs, RAG\n"
            "- Preferred Skills: Learning to Rank, LoRA/QLoRA, PEFT\n"
            "- Experience: 4+ years\n"
            "- Work Mode: In-office (Bangalore)\n\n"
            "Key Candidate Signals:\n"
            "1. Semantic profile alignment with JD tasks.\n"
            "2. Skill matching using recruiter-style synonym expansion.\n"
            "3. Career growth and stability (average job tenure > 2 years).\n"
            "4. Notice period availability (shorter notice preferred)."
        ),
        
        # Slide 4
        "How does your system retrieve, score, and rank candidates?": (
            "Ranking Methodology:\n"
            "- Retrieval (Stage 1): Precomputed BGE-Small-en-v1.5 embeddings for fast cosine similarity filtering of the top 500 candidates.\n"
            "- Scoring (Heuristics): Hybrid evaluation based on:\n"
            "  * 35% Semantic Similarity (BGE dense embeddings)\n"
            "  * 35% Skill Match (Must-have and preferred skills with synonym expansion)\n"
            "  * 20% Behavioral Fit (Work mode, notice period, and target salary)\n"
            "  * 10% Career Stability & Growth (Average job tenure and progressive title history)\n"
            "- Reranking (Stage 2): Reranking the top 500 candidates via ms-marco-MiniLM-L-6-v2 Cross-Encoder."
        ),
        
        # Slide 5
        "How are ranking decisions explained?": (
            "Explainability & Suspicious Profile Handling:\n"
            "- Explainability: Generates human-like recruiter justification citing exact current title/employer, total experience, matched must-have skills, average job tenure, and notice period.\n"
            "- Preventing Hallucinations: The generator reads facts directly from the pre-processed features (no LLM text generation is used in ranking, preventing any hallucination or variance).\n"
            "- Suspect Profiles: Instantly drops profiles to 0.0 score if they fail chronological checks (e.g. working at Krutrim in 2020) or show skill inflation (expert status with 0 months experience)."
        ),
        
        # Slide 6
        "What is the complete workflow": (
            "End-to-End Workflow:\n"
            "1. Input: Read docx/txt job description.\n"
            "2. Parsing: Extract role details, core skills, preferred skills, and experience criteria.\n"
            "3. Scoring: Compute hybrid score (semantic + skill + behavior + career) on precomputed database.\n"
            "4. Trap Filtering: Apply timeline and anomaly checks to filter fake candidates.\n"
            "5. Reranking: Perform Cross-Encoder reranking on the top 500 candidates.\n"
            "6. Output: Produce sorted CSV/XLSX ranking file with deterministic tie-breaking and recruiter reasoning."
        ),
        
        # Slide 7
        "System Architecture": (
            "System Architecture & Data Flow:\n\n"
            "1. Offline Precomputation:\n"
            "   - Extract candidate profile features (duration, stability, skills, work mode).\n"
            "   - Encode candidate profile summaries using BAAI/bge-small-en-v1.5 model.\n"
            "   - Save features to JSONL and embeddings to NumPy binary file.\n\n"
            "2. Online Pipeline (Executed on CPU under 35 seconds):\n"
            "   - Parse Job Description requirements.\n"
            "   - Query precomputed candidate embeddings to filter top candidates.\n"
            "   - Apply skill expansion, recruiter scoring heuristics, and honeypot checks.\n"
            "   - Pass top 500 candidates through ms-marco-MiniLM-L-6-v2 Cross-Encoder.\n"
            "   - Sort results, resolve tie-breaks, generate explainability reasoning, and save final top 100."
        ),
        
        # Slide 8
        "What results or insights": (
            "Results & Constraints Performance:\n"
            "- Computational Speed: Runs candidate scoring and Cross-Encoder reranking of the top 500 in 33.75 seconds on CPU (limit: 5 minutes).\n"
            "- Memory Constraint: Uses memory-efficient streaming generators to stay well within the 16GB RAM budget.\n"
            "- Validation: Fully validated by the official validation suite (exactly 100 candidates, sorted non-increasing scores, deterministic tie-breaking)."
        ),
        
        # Slide 9
        "What technologies, frameworks": (
            "Technology Stack & Rationale:\n"
            "- PyTorch & Sentence-Transformers: Selected for high-performance offline embedding generation and online Cross-Encoder reranking.\n"
            "- BAAI/bge-small-en-v1.5: SOTA lightweight embedding model chosen for excellent semantic representation and fast retrieval on CPU.\n"
            "- ms-marco-MiniLM-L-6-v2: Best-in-class lightweight Cross-Encoder for capturing detailed query-document interactions.\n"
            "- Pandas & NumPy: For efficient matrix manipulation and structured data saving.\n"
            "- Docker: For sandboxed, reproducible execution."
        ),
        
        # Slide 10
        "Github video etc": (
            "Submission Assets:\n"
            "- GitHub Repository: https://github.com/ravikrishna290/redrob-candidate-ranker\n"
            "- Docker Hub Image / Recipe: Build and run instructions included in README.md for 1-click sandbox testing.\n"
            "- Final Output File: team_redrob.xlsx (validated 100-candidate ranking list with recruiter justifications)."
        )
    }

    # Iterate through all slides and shapes
    for i, slide in enumerate(prs.slides, 1):
        print(f"Processing Slide {i}...")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Check for match in replacements
                matched_key = None
                for key in replacements:
                    if key in shape.text:
                        matched_key = key
                        break
                
                if matched_key:
                    print(f"  Replacing text in shape '{shape.name}': '{matched_key}'")
                    # Clear shape text frame and write replacement text
                    text_frame = shape.text_frame
                    text_frame.clear()
                    
                    # Add new paragraph and write content
                    p = text_frame.paragraphs[0]
                    p.text = replacements[matched_key]
                    
    # Save the populated presentation
    output_path = r"d:\Redrob AI hackathon\Idea Submission Template _ Redrob_Filled.pptx"
    prs.save(output_path)
    print(f"\nSuccessfully populated and saved PowerPoint file to: {output_path}")

if __name__ == "__main__":
    fill_slide_content()
