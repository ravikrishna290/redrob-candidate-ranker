import collections
import collections.abc
from pptx import Presentation

def main():
    prs = Presentation(r"d:\Redrob AI hackathon\Idea Submission Template _ Redrob.pptx")
    for i, slide in enumerate(prs.slides, 1):
        print(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(f"[{shape.name}]: {shape.text}")
        print()

if __name__ == "__main__":
    main()
